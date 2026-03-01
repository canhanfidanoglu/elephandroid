import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from .models import PlanReport


def _add_title_textbox(slide, text: str):
    """Add a title textbox at the top of a blank slide."""
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT


def _style_header_row(table, col_count: int):
    """Bold the first row of a table."""
    for col_idx in range(col_count):
        cell = table.cell(0, col_idx)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(12)


def build_pptx(report: PlanReport) -> bytes:
    """Build a PPTX presentation from a PlanReport and return as bytes."""
    prs = Presentation()

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = report.plan_name
    slide.placeholders[1].text = (
        f"Generated: {report.generated_at:%Y-%m-%d %H:%M} | "
        f"Overall: {report.overall_percentage:.0f}% complete"
    )

    # --- Slide 2: Task Summary ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    _add_title_textbox(slide, "Task Summary")

    not_started = report.total_tasks - report.completed_tasks - sum(
        b.in_progress for b in report.buckets
    )
    in_progress = sum(b.in_progress for b in report.buckets)

    rows = [
        ("Total Tasks", str(report.total_tasks)),
        ("Completed", str(report.completed_tasks)),
        ("In Progress", str(in_progress)),
        ("Not Started", str(not_started)),
    ]

    table_shape = slide.shapes.add_table(
        len(rows) + 1, 2, Inches(1.5), Inches(1.5), Inches(5), Inches(0.4 * (len(rows) + 1))
    )
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Count"
    _style_header_row(table, 2)
    for i, (metric, count) in enumerate(rows, start=1):
        table.cell(i, 0).text = metric
        table.cell(i, 1).text = count

    # --- Slide 3: Epic Progress ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    _add_title_textbox(slide, "Epic Progress")

    if report.epics:
        table_shape = slide.shapes.add_table(
            len(report.epics) + 1, 4,
            Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * (len(report.epics) + 1)),
        )
        table = table_shape.table
        headers = ["Epic", "Total", "Completed", "Percentage"]
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
        _style_header_row(table, 4)
        for i, epic in enumerate(report.epics, start=1):
            table.cell(i, 0).text = epic.name
            table.cell(i, 1).text = str(epic.total)
            table.cell(i, 2).text = str(epic.completed)
            table.cell(i, 3).text = f"{epic.percentage:.0f}%"

    # --- Slide 4: Bucket Breakdown ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    _add_title_textbox(slide, "Bucket Breakdown")

    if report.buckets:
        table_shape = slide.shapes.add_table(
            len(report.buckets) + 1, 5,
            Inches(0.3), Inches(1.5), Inches(9.4), Inches(0.4 * (len(report.buckets) + 1)),
        )
        table = table_shape.table
        headers = ["Bucket", "Total", "Completed", "In Progress", "Not Started"]
        for col_idx, header in enumerate(headers):
            table.cell(0, col_idx).text = header
        _style_header_row(table, 5)
        for i, bucket in enumerate(report.buckets, start=1):
            table.cell(i, 0).text = bucket.name
            table.cell(i, 1).text = str(bucket.total)
            table.cell(i, 2).text = str(bucket.completed)
            table.cell(i, 3).text = str(bucket.in_progress)
            table.cell(i, 4).text = str(bucket.not_started)

    # --- Slide 5: Timeline ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    _add_title_textbox(slide, "Task Timeline")
    txbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = "See Planner for detailed timeline."
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER

    # Serialize to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
